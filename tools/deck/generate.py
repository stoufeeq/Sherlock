"""Generate the Sherlock CTO pitch deck — UBS-branded .pptx.

Run:
    python -m venv .venv && source .venv/bin/activate
    pip install -r requirements.txt
    python generate.py
Output:
    ../../Sherlock_CTO_Pitch.pptx  (relative to this script)

Design notes:
- 16:9 (13.333 x 7.5 in).
- UBS Red #EC0016 as accent; black titles; grey body; lots of whitespace.
- Every slide has a top-left UBS text-mark — REPLACE with the official brand
  asset (click the group, Insert Picture, resize) before presenting.
- Speaker notes live in the notes pane for every slide.
- Fonts default to Arial (universal fallback). If UBS Sans / Frutiger UBS is
  installed on the presenting machine, change FONT_SANS below and regenerate.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# ---------- brand palette & typography ----------------------------------------
UBS_RED   = RGBColor(0xEC, 0x00, 0x16)
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLK  = RGBColor(0x17, 0x17, 0x17)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MID_GREY  = RGBColor(0x66, 0x66, 0x66)
LINE_GREY = RGBColor(0xD0, 0xD0, 0xD0)
BG_GREY   = RGBColor(0xF5, 0xF5, 0xF5)
CARD_BG   = RGBColor(0xFA, 0xFA, 0xFA)

FONT_SANS = "Arial"   # swap to "UBS Sans" or "Frutiger UBS" if installed
FONT_MONO = "Consolas"

OUT_PATH = Path(__file__).resolve().parents[2] / "Sherlock_CTO_Pitch.pptx"


# ---------- low-level helpers -------------------------------------------------

def _set_run(run, text, *, font=FONT_SANS, size=14, bold=False, italic=False, color=DARK_GREY):
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color


def _add_text(slide, left, top, width, height, text, *, font=FONT_SANS, size=14,
              bold=False, italic=False, color=DARK_GREY, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _set_run(run, text, font=font, size=size, bold=bold, italic=italic, color=color)
    return tb


def _add_bullets(slide, left, top, width, height, items, *, size=13, color=DARK_GREY,
                 bullet=True, line_spacing=1.25):
    """items = list of strings OR list of (head, body) tuples for head-bold bullets."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(2)
        if isinstance(item, tuple):
            head, body = item
            if bullet:
                r0 = p.add_run(); _set_run(r0, "— ", font=FONT_SANS, size=size, color=UBS_RED, bold=True)
            r1 = p.add_run(); _set_run(r1, head, font=FONT_SANS, size=size, color=NEAR_BLK, bold=True)
            r2 = p.add_run(); _set_run(r2, "  " + body, font=FONT_SANS, size=size, color=color)
        else:
            if bullet:
                r0 = p.add_run(); _set_run(r0, "— ", font=FONT_SANS, size=size, color=UBS_RED, bold=True)
            r = p.add_run()
            _set_run(r, item, font=FONT_SANS, size=size, color=color)
    return tb


def _add_rect(slide, left, top, width, height, *, fill=None, line=None, line_weight=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_weight)
    # remove shadow
    sp = shp.shadow
    sp.inherit = False
    return shp


def _add_connector(slide, x1, y1, x2, y2, *, color=MID_GREY, weight=1.5, arrow=True):
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    if arrow:
        # end arrow on
        ln = conn.line._get_or_add_ln()
        tailend = etree.SubElement(ln, qn("a:tailEnd"))
        tailend.set("type", "triangle")
        tailend.set("w", "med")
        tailend.set("h", "med")
    return conn


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------- slide chrome ------------------------------------------------------

def _header(slide, section_label="MAIN"):
    # UBS text-mark group (REPLACE with official logo)
    # three small red accent bars (approximation of the UBS keys visual cue)
    for i in range(3):
        y = Inches(0.22 + i * 0.11)
        _add_rect(slide, Inches(0.32), y, Inches(0.22), Inches(0.04), fill=UBS_RED)
    # UBS letters
    _add_text(slide, Inches(0.60), Inches(0.18), Inches(1.2), Inches(0.45),
              "UBS", font=FONT_SANS, size=22, bold=True, color=BLACK,
              anchor=MSO_ANCHOR.TOP)
    # section label + placeholder caption (right)
    _add_text(slide, Inches(10.6), Inches(0.25), Inches(2.5), Inches(0.3),
              section_label, font=FONT_SANS, size=9, color=MID_GREY,
              align=PP_ALIGN.RIGHT)


def _footer(slide, page, total, section="MAIN"):
    _add_rect(slide, Inches(0), Inches(7.2), Inches(13.333), Inches(0.02), fill=LINE_GREY)
    _add_text(slide, Inches(0.6), Inches(7.25), Inches(7), Inches(0.25),
              "Sherlock · Dependency Intelligence Platform · Internal · Draft",
              size=9, color=MID_GREY)
    _add_text(slide, Inches(10.3), Inches(7.25), Inches(2.5), Inches(0.25),
              f"{section} · {page} / {total}", size=9, color=MID_GREY,
              align=PP_ALIGN.RIGHT)


def _title_block(slide, title, eyebrow=None, subtitle=None):
    if eyebrow:
        _add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.3),
                  eyebrow.upper(), size=10, bold=True, color=UBS_RED)
    _add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.9),
              title, size=30, bold=True, color=NEAR_BLK)
    # red keyline
    _add_rect(slide, Inches(0.6), Inches(2.05), Inches(0.9), Inches(0.06), fill=UBS_RED)
    if subtitle:
        _add_text(slide, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
                  subtitle, size=13, color=MID_GREY)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------- slide builders ----------------------------------------------------

def build_title_slide(prs, idx, total):
    s = _blank_slide(prs)
    # black background
    _add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill=NEAR_BLK)
    # accent bar
    _add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.15), fill=UBS_RED)

    # UBS mark in white on dark
    for i in range(3):
        y = Inches(0.55 + i * 0.13)
        _add_rect(s, Inches(0.60), y, Inches(0.26), Inches(0.05), fill=UBS_RED)
    _add_text(s, Inches(0.95), Inches(0.50), Inches(2), Inches(0.5),
              "UBS", size=26, bold=True, color=WHITE)

    # title block
    _add_text(s, Inches(0.6), Inches(2.6), Inches(12), Inches(0.5),
              "CTO OFFICE BRIEFING · PROPOSAL", size=11, bold=True, color=UBS_RED)
    _add_text(s, Inches(0.6), Inches(3.0), Inches(12), Inches(1.6),
              "Sherlock", size=64, bold=True, color=WHITE)
    _add_text(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.7),
              "Dependency Intelligence Platform", size=24, color=WHITE)

    # red keyline
    _add_rect(s, Inches(0.6), Inches(5.1), Inches(1.6), Inches(0.06), fill=UBS_RED)

    _add_text(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.5),
              "Cross-application change impact — surfaced at commit time.",
              size=16, italic=True, color=WHITE)

    _add_text(s, Inches(0.6), Inches(6.4), Inches(12), Inches(0.35),
              "Presented to the CTO Office  ·  Internal briefing  ·  Draft",
              size=10, color=LINE_GREY)

    _notes(s, "Opening line: 'Over the last year we have lost engineering time and "
              "incurred customer-visible incidents every time a change in one application "
              "silently broke another. Sherlock closes that gap — at commit time, not in production.' "
              "This is a 3-month pilot ask and a platform-embed request; I'll get to both "
              "on slide 14.")
    return s


def build_exec_summary(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "What we need in 60 seconds",
                 eyebrow="Executive summary")

    # three cards
    x0, y0, w, h, gap = Inches(0.6), Inches(2.8), Inches(4.0), Inches(3.6), Inches(0.15)
    cards = [
        ("Problem",
         "Code changes in one app silently break others.",
         "Dependencies are invisible until something fails in production. "
         "We detect the break hours-to-days after merge, not at merge time. "
         "Every domain has recurring examples; nobody has end-to-end visibility."),
        ("Solution",
         "Sherlock — the live dependency graph of UBS.",
         "Auto-maps every application's upstream/downstream surface from source "
         "code: REST, events, DB, shared files, libraries. Flags potential "
         "breakage on every MR, with the right team tagged and on-call channel named."),
        ("Ask",
         "3-month funded pilot · 2–3 FTEs · platform embed.",
         "One pilot domain, measurable MTTD reduction, ready-to-scale platform. "
         "Approval to embed in the existing GitLab CI stack and executive sponsor "
         "to unblock cross-BU alignment."),
    ]
    for i, (head, lead, body) in enumerate(cards):
        x = x0 + (w + gap) * i
        _add_rect(s, x, y0, w, h, fill=CARD_BG, line=LINE_GREY)
        _add_rect(s, x, y0, w, Inches(0.08), fill=UBS_RED)
        _add_text(s, x + Inches(0.3), y0 + Inches(0.25), w - Inches(0.6), Inches(0.35),
                  head.upper(), size=10, bold=True, color=UBS_RED)
        _add_text(s, x + Inches(0.3), y0 + Inches(0.7), w - Inches(0.6), Inches(1.1),
                  lead, size=16, bold=True, color=NEAR_BLK)
        _add_text(s, x + Inches(0.3), y0 + Inches(1.75), w - Inches(0.6), h - Inches(1.9),
                  body, size=12, color=DARK_GREY)

    _notes(s, "If the CTO needs to leave after one slide, this is the one they'd take with "
              "them. Lead with the ask in plain language; the rest of the deck defends it. "
              "Scale anchor: '10,000 engineers, 17 domains, thousands of repos — nobody sees "
              "the whole picture.'")
    return s


def build_problem(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "How cross-application breakage happens — today",
                 eyebrow="The problem",
                 subtitle="Recurring patterns across 17 domains, 10,000+ engineers, 1000s of repos. "
                          "Composite illustrative examples; every domain has their own version.")

    patterns = [
        ("1. Contract drift",
         "Team A renames a response field in its OpenAPI. Teams B & C find out the next morning "
         "when their nightly integration tests fail, or — worse — when prod errors spike.",
         "Typical MTTD: 1–3 days · MTTR: 8–24 hours"),
        ("2. Silent DB coupling",
         "Team A owns a Postgres table. Team B's analytics app reads it directly for a "
         "regulatory report. An 'internal' column rename lands on main — the report breaks "
         "two days later. Neither team knew the other was there.",
         "Typical MTTD: 2–5 days · MTTR: 12–48 hours"),
        ("3. Legacy file feeds",
         "A modern service stops writing a nightly batch file. A core-banking COBOL job fails "
         "overnight. Ops notices the morning after when reports don't land.",
         "Typical MTTD: 1 day (overnight batch) · MTTR: 12–36 hours"),
    ]
    y = Inches(3.1)
    for i, (h, body, cost) in enumerate(patterns):
        row_y = y + Inches(i * 1.3)
        _add_rect(s, Inches(0.6), row_y, Inches(0.1), Inches(1.1), fill=UBS_RED)
        _add_text(s, Inches(0.9), row_y, Inches(5.0), Inches(0.35), h,
                  size=14, bold=True, color=NEAR_BLK)
        _add_text(s, Inches(0.9), row_y + Inches(0.35), Inches(8.7), Inches(0.65),
                  body, size=11, color=DARK_GREY)
        _add_text(s, Inches(9.8), row_y + Inches(0.35), Inches(3.3), Inches(0.65),
                  cost, size=10, italic=True, color=UBS_RED, align=PP_ALIGN.RIGHT)

    _notes(s, "Say: 'Every one of these patterns is already happening somewhere in our estate "
              "every week.' Do NOT attach to a specific incident unless you've cleared one — the "
              "composite framing is more defensible. Invite them to recall one from their own "
              "domain if they'd like to share.")
    return s


def build_gap(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "The visibility gap nobody owns",
                 eyebrow="Why existing tools don't close this",
                 subtitle="Existing tools each solve part of the picture. None closes the loop "
                          "between code-change and cross-app impact, at MR time.")

    rows = [
        ("CMDB / service catalog",  "Good for ownership. Stale topology within weeks — "
                                    "it's manually maintained."),
        ("APM (Dynatrace)",         "Shows runtime calls that actually executed. Silent on "
                                    "files, batch, dormant paths, legacy feeds."),
        ("Code search / Sourcegraph","Finds references. Doesn't understand contracts, directionality, "
                                    "or cross-app impact."),
        ("Architecture diagrams",   "Drawn once. Never updated. Confidence decays to zero."),
        ("Backstage / dev portals", "Great for ownership + docs. Doesn't parse contracts or "
                                    "diff them across MRs."),
    ]
    # header row
    y = Inches(3.0)
    _add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.45), fill=NEAR_BLK)
    _add_text(s, Inches(0.8), y + Inches(0.08), Inches(4.0), Inches(0.3), "TOOL",
              size=10, bold=True, color=WHITE)
    _add_text(s, Inches(5.0), y + Inches(0.08), Inches(8.0), Inches(0.3), "GAP",
              size=10, bold=True, color=WHITE)
    for i, (tool, gap) in enumerate(rows):
        ry = y + Inches(0.5 + i * 0.55)
        _add_rect(s, Inches(0.6), ry, Inches(12.1), Inches(0.5),
                  fill=BG_GREY if i % 2 == 0 else WHITE, line=LINE_GREY)
        _add_text(s, Inches(0.8), ry + Inches(0.1), Inches(4.0), Inches(0.35),
                  tool, size=12, bold=True, color=NEAR_BLK)
        _add_text(s, Inches(5.0), ry + Inches(0.1), Inches(7.9), Inches(0.35),
                  gap, size=11, color=DARK_GREY)

    _add_text(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.45),
              "No single tool connects CODE  →  CROSS-APP IMPACT  →  MR-TIME FEEDBACK. That's the gap Sherlock closes.",
              size=13, bold=True, italic=True, color=UBS_RED)

    _notes(s, "Expect the CTO to ask how Sherlock relates to each of these. Bridge answer: "
              "'Sherlock sits on top of the code — it complements CMDB (which it pulls "
              "ownership from), complements APM (which it will later cross-validate against), "
              "and feeds Backstage (it can expose its graph as a plugin).'")
    return s


def build_intro_sherlock(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "Introducing Sherlock",
                 eyebrow="The idea")

    _add_text(s, Inches(0.6), Inches(2.5), Inches(12.1), Inches(0.6),
              "The enterprise's live dependency graph — continuously derived from source code.",
              size=18, italic=True, color=NEAR_BLK)

    # three pillars
    pillars = [
        ("SEE", "The real topology across all 17 domains.",
         "REST, events, DB, shared files, libraries — one graph."),
        ("PREDICT", "The blast radius of every MR, before it merges.",
         "Automatic breaking-change detection. Source app & affected apps named with on-call channels."),
        ("NOTIFY", "The right teams — automatically.",
         "MR comment for the author. Sticky impact issue in every affected repo. "
         "Auto-closes when the fix lands."),
    ]
    y0 = Inches(3.5)
    w = Inches(4.0); gap = Inches(0.15); x0 = Inches(0.6)
    for i, (head, lead, body) in enumerate(pillars):
        x = x0 + (w + gap) * i
        _add_rect(s, x, y0, w, Inches(2.8), fill=WHITE, line=LINE_GREY)
        _add_rect(s, x, y0, Inches(0.08), Inches(2.8), fill=UBS_RED)
        _add_text(s, x + Inches(0.3), y0 + Inches(0.2), w - Inches(0.6), Inches(0.45),
                  head, size=14, bold=True, color=UBS_RED)
        _add_text(s, x + Inches(0.3), y0 + Inches(0.7), w - Inches(0.6), Inches(0.7),
                  lead, size=14, bold=True, color=NEAR_BLK)
        _add_text(s, x + Inches(0.3), y0 + Inches(1.5), w - Inches(0.6), Inches(1.2),
                  body, size=11, color=DARK_GREY)

    _add_text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.35),
              "No agents, no code changes in consuming apps. Reads code. Writes to GitLab.",
              size=11, italic=True, color=MID_GREY)

    _notes(s, "Three-word pitch: SEE, PREDICT, NOTIFY. "
              "Stress the zero-intrusion point — teams don't have to instrument anything; "
              "Sherlock works from the code that's already in GitLab.")
    return s


def build_architecture(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "How Sherlock works", eyebrow="Architecture",
                 subtitle="Webhook-driven. Internal only. Read-only from applications' perspective.")

    # blocks positions (inches)
    # GitLab → Ingest → Analyzers → Graph → Impact Engine → Notifier → GitLab (feedback loop)
    # Canvas reads Graph
    def box(x, y, w, h, head, sub="", accent=UBS_RED, fill=WHITE):
        _add_rect(s, Inches(x), Inches(y), Inches(w), Inches(h), fill=fill, line=LINE_GREY)
        _add_rect(s, Inches(x), Inches(y), Inches(w), Inches(0.08), fill=accent)
        _add_text(s, Inches(x + 0.15), Inches(y + 0.18), Inches(w - 0.3), Inches(0.35),
                  head, size=11, bold=True, color=NEAR_BLK)
        if sub:
            _add_text(s, Inches(x + 0.15), Inches(y + 0.52), Inches(w - 0.3), Inches(h - 0.6),
                      sub, size=9, color=MID_GREY)

    row_y = 3.1
    # row 1: GitLab <-> Sherlock
    box(0.6, row_y, 2.2, 1.1, "GitLab", "Webhooks: push, MR\nSource of truth", accent=BLACK, fill=BG_GREY)
    box(3.3, row_y, 2.2, 1.1, "Ingest (FastAPI)", "Webhook receiver\nManual /scan")
    box(6.0, row_y, 2.2, 1.1, "Analyzers", "Java · Python\nJS/TS · COBOL · YAML")
    box(8.7, row_y, 2.2, 1.1, "Graph DB (Neo4j)", "Versioned by commit\nApp / Endpoint / Topic\nTable / File / Library")
    box(11.1, row_y, 1.6, 1.1, "Canvas UI", "Cytoscape.js\nApp + Contract views")

    # arrows left to right
    for x1, x2 in [(2.8, 3.3), (5.5, 6.0), (8.2, 8.7), (10.9, 11.1)]:
        _add_connector(s, Inches(x1), Inches(row_y + 0.55), Inches(x2), Inches(row_y + 0.55))

    # row 2: Impact + Notifier feeding back to GitLab
    row2_y = 5.2
    box(3.3, row2_y, 2.8, 1.1, "Impact Engine", "Diff AnalysisResult\nResolve affected apps", accent=UBS_RED)
    box(6.4, row2_y, 2.8, 1.1, "Notifier", "MR comments\nSticky impact issues", accent=UBS_RED)
    box(9.5, row2_y, 2.2, 1.1, "CMDB", "Team · tier · on-call\nRead-only", accent=BLACK, fill=BG_GREY)

    # arrows from graph → impact → notifier → gitlab (back up)
    _add_connector(s, Inches(9.7), Inches(row_y + 1.1), Inches(4.7), Inches(row2_y))
    _add_connector(s, Inches(6.1), Inches(row2_y + 0.55), Inches(6.4), Inches(row2_y + 0.55))
    _add_connector(s, Inches(7.8), Inches(row2_y), Inches(1.6), Inches(row_y + 1.1))
    _add_connector(s, Inches(9.5), Inches(row2_y + 0.55), Inches(9.4), Inches(row2_y + 0.55),
                   arrow=False)
    _add_connector(s, Inches(9.5), Inches(row2_y + 0.2), Inches(9.5), Inches(row_y + 1.1), arrow=False)

    _add_text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.35),
              "All in-tenant. Source code never leaves the UBS perimeter. "
              "No PII in the graph — code, not runtime data.",
              size=10, italic=True, color=MID_GREY)

    _notes(s, "Walk the diagram left-to-right first (how we build the graph), then the "
              "feedback loop (how impact gets surfaced). Emphasize: CMDB is read-only input "
              "— we don't try to replace it. APM integration is roadmap (H4).")
    return s


def build_coverage(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "What Sherlock detects today",
                 eyebrow="Coverage matrix",
                 subtitle="Multi-language, multi-surface detection. COBOL coverage is a key "
                          "differentiator for our core-banking estate.")

    headers = ["COUPLING TYPE", "SIGNAL SOURCE", "LANGUAGES", "BREAK KINDS DETECTED"]
    rows = [
        ("REST", "OpenAPI + HTTP client code (regex + tree-sitter AST)", "Java · Python · JS/TS",
         "Endpoint removed · Required-field shape changed · Optional-field added (info) · Deprecated removal (info)"),
        ("Events (Kafka)", "AsyncAPI + publish/subscribe code", "Java · Python",
         "Topic publish removed · Required-field payload changed · Optional-field added (info)"),
        ("Database", "Flyway migrations + SQL in code", "Java · Python · COBOL",
         "Table write removed · Schema ownership released"),
        ("Shared files", "Path literals + COBOL SELECT/ASSIGN", "Java · Python · JS · COBOL",
         "File feed write removed"),
        ("Shared libraries", "pom.xml · requirements.txt · package.json", "Java · Python · Node",
         "Published library retired / renamed"),
    ]
    x_cols = [0.6, 3.1, 6.3, 8.5]
    w_cols = [2.5, 3.2, 2.2, 4.2]
    y = Inches(3.0)

    _add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.45), fill=NEAR_BLK)
    for i, h in enumerate(headers):
        _add_text(s, Inches(x_cols[i] + 0.15), y + Inches(0.1),
                  Inches(w_cols[i]), Inches(0.3),
                  h, size=9, bold=True, color=WHITE)

    for i, row in enumerate(rows):
        ry = y + Inches(0.5 + i * 0.65)
        _add_rect(s, Inches(0.6), ry, Inches(12.1), Inches(0.6),
                  fill=BG_GREY if i % 2 == 0 else WHITE, line=LINE_GREY)
        for j, text in enumerate(row):
            is_cobol = "COBOL" in text and j != 3
            color = UBS_RED if is_cobol else NEAR_BLK if j == 0 else DARK_GREY
            _add_text(s, Inches(x_cols[j] + 0.15), ry + Inches(0.15),
                      Inches(w_cols[j]), Inches(0.5), text,
                      size=10, bold=(j == 0 or is_cobol), color=color)

    _add_text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.35),
              "Required-field changes are flagged BREAKING; optional-field additions land as INFO — "
              "the diff engine separates contract evolution from contract regression.",
              size=10, italic=True, color=MID_GREY)

    _notes(s, "Highlight three things: (1) COBOL line — few tools in this space speak COBOL. "
              "(2) tree-sitter — Java/Python/JS now extracted at AST level, catching URI-builder lambdas + "
              "variable-resolved URLs the regex extractor misses. (3) precision — required-vs-optional "
              "decomposition means adding a new optional response field doesn't fire a breaking-change alert.")
    return s


def build_mr_moment(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "Impact surfaced at MR time", eyebrow="The feedback loop",
                 subtitle="When a team opens an MR that would break a downstream app, "
                          "Sherlock posts a structured comment on the MR — with affected teams "
                          "and on-call channels named.")

    # Left: mock MR comment — now showing the cross-platform callout banner
    left_x = Inches(0.6); left_y = Inches(2.85); left_w = Inches(6.5); left_h = Inches(4.2)
    _add_rect(s, left_x, left_y, left_w, left_h, fill=CARD_BG, line=LINE_GREY)
    _add_rect(s, left_x, left_y, left_w, Inches(0.08), fill=UBS_RED)

    _add_text(s, left_x + Inches(0.3), left_y + Inches(0.18), left_w - Inches(0.6), Inches(0.3),
              "SHERLOCK BOT  ·  banking/account-service !12", size=9, bold=True, color=UBS_RED)
    _add_text(s, left_x + Inches(0.3), left_y + Inches(0.5), left_w - Inches(0.6), Inches(0.45),
              "🔎 Sherlock Impact Analysis", size=15, bold=True, color=NEAR_BLK)
    _add_text(s, left_x + Inches(0.3), left_y + Inches(0.95), left_w - Inches(0.6), Inches(0.3),
              "Source: account-service · platform on-prem · target main",
              size=10, color=MID_GREY, font=FONT_MONO)

    _add_text(s, left_x + Inches(0.3), left_y + Inches(1.3), left_w - Inches(0.6), Inches(0.3),
              "⚠️  1 breaking change detected",
              size=11, bold=True, color=UBS_RED)

    # Cross-platform callout — UBS-red banner inside the comment
    callout_y = left_y + Inches(1.65)
    _add_rect(s, left_x + Inches(0.25), callout_y, left_w - Inches(0.5), Inches(0.55),
              fill=RGBColor(0xFE, 0xE2, 0xE2), line=UBS_RED)
    _add_text(s, left_x + Inches(0.4), callout_y + Inches(0.06), left_w - Inches(0.7), Inches(0.5),
              "🚨  Cross-platform impact — originates on  on-prem  and affects 3 app(s) on  azure.",
              size=10, bold=True, color=UBS_RED)

    _add_text(s, left_x + Inches(0.3), left_y + Inches(2.35), left_w - Inches(0.6), Inches(0.3),
              "🔌 Removed REST endpoint",
              size=11, bold=True, color=NEAR_BLK)
    _add_text(s, left_x + Inches(0.3), left_y + Inches(2.65), left_w - Inches(0.6), Inches(0.35),
              "GET /accounts/{*}", size=10, color=DARK_GREY, font=FONT_MONO)

    _add_text(s, left_x + Inches(0.3), left_y + Inches(3.05), left_w - Inches(0.6), Inches(0.3),
              "Directly affected (4):", size=10, bold=True, color=NEAR_BLK)
    _add_text(s, left_x + Inches(0.3), left_y + Inches(3.35), left_w - Inches(0.6), Inches(0.35),
              "• mobile-bff      mobile-experience  T1  azure  🚨",
              size=9, color=DARK_GREY, font=FONT_MONO)
    _add_text(s, left_x + Inches(0.3), left_y + Inches(3.55), left_w - Inches(0.6), Inches(0.35),
              "• web-portal-bff  web-experience     T1  azure  🚨",
              size=9, color=DARK_GREY, font=FONT_MONO)
    _add_text(s, left_x + Inches(0.3), left_y + Inches(3.75), left_w - Inches(0.6), Inches(0.35),
              "• fraud-detection risk-engineering   T1  azure  🚨",
              size=9, color=DARK_GREY, font=FONT_MONO)
    _add_text(s, left_x + Inches(0.3), left_y + Inches(3.95), left_w - Inches(0.6), Inches(0.35),
              "• transaction-svc payments-platform  T1  on-prem",
              size=9, color=DARK_GREY, font=FONT_MONO)

    # right: explanatory bullets
    right_x = Inches(7.4); right_y = Inches(2.85)
    _add_text(s, right_x, right_y, Inches(5.3), Inches(0.4),
              "What this replaces",
              size=12, bold=True, color=UBS_RED)
    _add_bullets(s, right_x, right_y + Inches(0.5), Inches(5.3), Inches(3.4), [
        ("Slack triage.", "No more 'who touched X yesterday?' threads across 17 domains."),
        ("Blame-after-the-fact.", "Affected teams see the impact before the author hits Merge."),
        ("Gated releases.", "Soft signal — inform, don't block. Keeps velocity."),
        ("Knowledge silos.", "On-call channel is named in the comment, not discovered."),
        ("Cross-cloud surprise.", "🚨 banner explicitly flags Azure↔on-prem boundary crossings — "
                                  "the costliest type of break to debug."),
    ], size=12, line_spacing=1.35)

    _notes(s, "This is the 'aha' slide. If possible, open http://localhost:8001 in a browser "
              "and show the live canvas during this moment. If not, the mock here is accurate "
              "to what a real MR comment looks like — cross-platform callout banner included.")
    return s


def build_canvas(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "The canvas — dependency visibility for every team",
                 eyebrow="The enterprise topology",
                 subtitle="App view · Contract view · Team filter · Downstream/Upstream impact overlay")

    # placeholder for canvas screenshot
    ph_x, ph_y, ph_w, ph_h = Inches(0.6), Inches(3.0), Inches(7.5), Inches(3.8)
    _add_rect(s, ph_x, ph_y, ph_w, ph_h, fill=BG_GREY, line=LINE_GREY)
    _add_text(s, ph_x, ph_y + Inches(1.7), ph_w, Inches(0.4),
              "[ insert canvas screenshot — http://localhost:8001/ui ]",
              size=12, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)
    _add_text(s, ph_x, ph_y + Inches(2.2), ph_w, Inches(0.4),
              "App view · rankDir: BT · upstream on top, downstream on bottom",
              size=10, color=MID_GREY, align=PP_ALIGN.CENTER)

    # right: capabilities
    right_x = Inches(8.4); right_y = Inches(3.0)
    _add_text(s, right_x, right_y, Inches(4.3), Inches(0.4),
              "Answered at a glance", size=12, bold=True, color=UBS_RED)
    _add_bullets(s, right_x, right_y + Inches(0.5), Inches(4.3), Inches(3.2), [
        "Who am I coupled to (upstream & downstream)?",
        "What breaks if my change ships?",
        "Who owns each endpoint, topic, table, feed?",
        "Which teams inherit the blast radius?",
        "Where are the cross-domain edges?",
    ], size=12, line_spacing=1.4)

    _notes(s, "Screenshot replacement: take a PNG of the canvas in App view, selected on "
              "transaction-service with Downstream overlay enabled. Crop to 16:9.")
    return s


def build_sticky_lifecycle(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "Nothing falls through the cracks",
                 eyebrow="Sticky impact tags",
                 subtitle="Every affected repo gets an auto-opened issue. It closes itself "
                          "when the underlying break is resolved.")

    # flow
    steps = [
        ("MR opens", "Break detected in source repo.\nMR comment posted.", UBS_RED),
        ("Issue opens", "Auto-opened in each affected repo with label impact::pending.\n"
                        "Linked to source MR.", UBS_RED),
        ("Team responds", "Downstream team sees in their queue.\n"
                          "Coordinates fix or accepts impact.", BLACK),
        ("Fix lands", "Rescan on main. Break no longer reproducible.\n"
                      "Issue auto-closed with label impact::fixed and resolution note.", RGBColor(0x10, 0xB9, 0x81)),
    ]
    y = Inches(3.1); w = Inches(2.9); gap = Inches(0.3); x0 = Inches(0.6); h = Inches(2.5)
    for i, (head, body, accent) in enumerate(steps):
        x = x0 + (w + gap) * i
        _add_rect(s, x, y, w, h, fill=WHITE, line=LINE_GREY)
        _add_rect(s, x, y, w, Inches(0.08), fill=accent)
        _add_text(s, x + Inches(0.2), y + Inches(0.25), w - Inches(0.4), Inches(0.35),
                  f"STEP {i+1}", size=9, bold=True, color=accent)
        _add_text(s, x + Inches(0.2), y + Inches(0.6), w - Inches(0.4), Inches(0.5),
                  head, size=15, bold=True, color=NEAR_BLK)
        _add_text(s, x + Inches(0.2), y + Inches(1.2), w - Inches(0.4), h - Inches(1.3),
                  body, size=11, color=DARK_GREY)
        if i < 3:
            _add_connector(s, x + w + Inches(0.02), y + Inches(1.2),
                           x + w + gap - Inches(0.02), y + Inches(1.2))

    _add_text(s, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.45),
              "SLA becomes measurable:  mean time to clear impact tag, per team.",
              size=13, bold=True, italic=True, color=UBS_RED)

    _notes(s, "This is the slide that sells 'inform, don't block'. The system gives teams "
              "the data; responsibility stays with the team. SLA dashboard comes in H3.")
    return s


def build_hybrid_platform(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "Hybrid estate, one graph",
                 eyebrow="Cross-platform awareness",
                 subtitle="Same GitLab. One canvas. Every app tagged with its deployment "
                          "platform — and every Azure ↔ on-prem boundary crossing called out.")

    AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
    SLATE      = RGBColor(0x6B, 0x72, 0x80)
    AZURE_BG   = RGBColor(0xE6, 0xF2, 0xFB)
    ON_PREM_BG = RGBColor(0xF1, 0xF3, 0xF5)

    # ---- Two-zone diagram (left ~7.6") -------------------------------------
    diag_x = Inches(0.6); diag_y = Inches(2.85); diag_w = Inches(7.6); diag_h = Inches(4.0)

    half_w = (diag_w - Inches(0.2)) / 2
    # Azure zone
    _add_rect(s, diag_x, diag_y, half_w, diag_h, fill=AZURE_BG, line=AZURE_BLUE)
    _add_rect(s, diag_x, diag_y, half_w, Inches(0.4), fill=AZURE_BLUE)
    _add_text(s, diag_x + Inches(0.2), diag_y + Inches(0.06),
              half_w - Inches(0.4), Inches(0.3),
              "AZURE", size=11, bold=True, color=WHITE)

    # On-prem zone
    onprem_x = diag_x + half_w + Inches(0.2)
    _add_rect(s, onprem_x, diag_y, half_w, diag_h, fill=ON_PREM_BG, line=SLATE)
    _add_rect(s, onprem_x, diag_y, half_w, Inches(0.4), fill=SLATE)
    _add_text(s, onprem_x + Inches(0.2), diag_y + Inches(0.06),
              half_w - Inches(0.4), Inches(0.3),
              "ON-PREM", size=11, bold=True, color=WHITE)

    # Application bubbles — coloured rings reflect the platform tag in CMDB
    def _app_node(left, top, label, ring):
        # ring rectangle (3pt outline)
        _add_rect(s, left, top, Inches(1.6), Inches(0.55),
                  fill=WHITE, line=ring, line_weight=2.5)
        _add_text(s, left, top + Inches(0.13), Inches(1.6), Inches(0.3),
                  label, size=10, bold=True, color=NEAR_BLK, align=PP_ALIGN.CENTER)

    # Azure-side apps (4)
    azure_inner_x = diag_x + Inches(0.3)
    _app_node(azure_inner_x,                 diag_y + Inches(0.7), "mobile-bff",       AZURE_BLUE)
    _app_node(azure_inner_x + Inches(1.8),   diag_y + Inches(0.7), "web-portal-bff",   AZURE_BLUE)
    _app_node(azure_inner_x,                 diag_y + Inches(1.55), "fraud-detection", AZURE_BLUE)
    _app_node(azure_inner_x + Inches(1.8),   diag_y + Inches(1.55), "analytics-svc",   AZURE_BLUE)

    # On-prem-side apps (5)
    op_inner_x = onprem_x + Inches(0.3)
    _app_node(op_inner_x,                    diag_y + Inches(0.7),  "account-service",     SLATE)
    _app_node(op_inner_x + Inches(1.8),      diag_y + Inches(0.7),  "transaction-service", SLATE)
    _app_node(op_inner_x,                    diag_y + Inches(1.55), "customer-service",    SLATE)
    _app_node(op_inner_x + Inches(1.8),      diag_y + Inches(1.55), "notification-svc",    SLATE)
    _app_node(op_inner_x + Inches(0.9),      diag_y + Inches(2.4),  "legacy-ledger",       SLATE)

    # Boundary-crossing edges — drawn in UBS red to flag cross-platform calls.
    # (account-service /accounts/{*} ← mobile-bff, web-portal-bff, fraud-detection)
    src_x = op_inner_x + Inches(0.8)
    src_y = diag_y + Inches(0.7) + Inches(0.275)
    for tgt in [(azure_inner_x + Inches(0.8), diag_y + Inches(0.7) + Inches(0.275)),       # mobile-bff
                (azure_inner_x + Inches(2.6), diag_y + Inches(0.7) + Inches(0.275)),       # web-portal-bff
                (azure_inner_x + Inches(0.8), diag_y + Inches(1.55) + Inches(0.275))]:     # fraud-detection
        _add_connector(s, src_x, src_y, tgt[0], tgt[1], color=UBS_RED, weight=1.75, arrow=True)

    # 🚨 callout label sitting between the two zones
    _add_rect(s, diag_x + half_w - Inches(0.4), diag_y + diag_h - Inches(0.55),
              Inches(1.0), Inches(0.4), fill=UBS_RED)
    _add_text(s, diag_x + half_w - Inches(0.4), diag_y + diag_h - Inches(0.49),
              Inches(1.0), Inches(0.3),
              "🚨 cross-boundary",
              size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # caption under the diagram
    _add_text(s, diag_x, Inches(6.95), diag_w, Inches(0.25),
              "Live edges from the POC: account-service (on-prem) is consumed by 3 azure-side apps.",
              size=9, italic=True, color=MID_GREY)

    # ---- Right rail — capability bullets ----------------------------------
    rx = Inches(8.4); ry = Inches(2.85); rw = Inches(4.4)
    _add_text(s, rx, ry, rw, Inches(0.4),
              "What's tagged, surfaced, and labelled",
              size=12, bold=True, color=UBS_RED)
    _add_bullets(s, rx, ry + Inches(0.5), rw, Inches(4.0), [
        ("CMDB drives it. ",     "Each app's `platform` field (azure / on-prem / library) "
                                  "merges into the graph on every scan."),
        ("Coloured rings. ",     "Azure apps ring blue, on-prem ring slate — instantly "
                                  "visible on the canvas."),
        ("Platform filter. ",    "Toolbar dropdown filters the canvas to one platform."),
        ("MR comment banner. ",  "🚨 Cross-platform impact when source platform ≠ "
                                  "affected platform — placed at the top of the comment."),
        ("Sticky-issue label. ", "`impact::cross-platform` (UBS red) auto-applied to "
                                  "downstream issues that span the boundary."),
        ("Same workflow. ",      "One GitLab, one webhook flow, one canvas — no separate "
                                  "tooling per platform."),
    ], size=11, line_spacing=1.35)

    _notes(s, "This slide answers the obvious follow-up to slide 8: 'How does this work in our "
              "hybrid estate?' Anchor: 'Sherlock doesn't care where you deploy. CMDB tells us "
              "the platform, the graph tracks the edges, and any time a break crosses the "
              "Azure ↔ on-prem boundary we surface it explicitly — at MR time and on the "
              "downstream issue. Cross-cloud incidents are the costliest to debug; this is "
              "the cheap end of catching them.'")
    return s


def build_autodoc(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(
        s,
        "Every app gets a self-maintaining README",
        eyebrow="Auto-documentation",
        subtitle="One-click per app · draft MR in THAT app's own repo · reviewed by THAT app's own team.",
    )

    # Diff-style preview card — real Gemini output from the legacy-ledger demo MR.
    # Uses a pale-green "added" strip + monospace body to evoke a GitLab diff pane.
    card_x = Inches(0.6); card_y = Inches(2.85); card_w = Inches(8.3); card_h = Inches(4.0)
    _add_rect(s, card_x, card_y, card_w, card_h, fill=WHITE, line=LINE_GREY)

    # Title bar of the MR preview
    _add_rect(s, card_x, card_y, card_w, Inches(0.45), fill=NEAR_BLK)
    _add_text(s, card_x + Inches(0.2), card_y + Inches(0.08), Inches(4.5), Inches(0.3),
              "📄  banking/legacy-ledger  ·  README.md",
              size=10, bold=True, color=WHITE, font=FONT_MONO)
    _add_text(s, card_x + Inches(4.9), card_y + Inches(0.08), Inches(3.3), Inches(0.3),
              "MR  sherlock::autodoc  ·  Draft",
              size=9, color=RGBColor(0xFC, 0xA5, 0xA5), align=PP_ALIGN.RIGHT)

    # Diff gutter (soft green stripe indicating "added content")
    gutter = RGBColor(0xD1, 0xFA, 0xE5)
    _add_rect(s, card_x, card_y + Inches(0.45), Inches(0.08), card_h - Inches(0.45), fill=gutter)

    # Content area
    cx = card_x + Inches(0.35); cy = card_y + Inches(0.7); cw = card_w - Inches(0.6)

    _add_text(s, cx, cy, cw, Inches(0.3),
              "<!-- sherlock:autodoc-start -->",
              size=9, color=MID_GREY, font=FONT_MONO)

    _add_text(s, cx, cy + Inches(0.35), cw, Inches(0.4),
              "## 🔎 Auto-generated by Sherlock",
              size=15, bold=True, color=NEAR_BLK)

    # The actual Gemini-written Purpose, quoted verbatim from the legacy-ledger MR
    purpose_card_x = cx
    purpose_card_y = cy + Inches(0.8)
    purpose_card_h = Inches(0.7)
    _add_rect(s, purpose_card_x, purpose_card_y, cw, purpose_card_h,
              fill=RGBColor(0xEE, 0xF2, 0xFF), line=RGBColor(0xC7, 0xD2, 0xFE))
    _add_text(s, purpose_card_x + Inches(0.2), purpose_card_y + Inches(0.08),
              cw - Inches(0.4), Inches(0.3),
              "PURPOSE  ·  written by Gemini 2.5 Flash",
              size=9, bold=True, color=RGBColor(0x37, 0x30, 0xA3))
    _add_text(s, purpose_card_x + Inches(0.2), purpose_card_y + Inches(0.33),
              cw - Inches(0.4), Inches(0.35),
              "\"The legacy-ledger application updates ledger balances and generates a ledger report.\"",
              size=12, italic=True, color=NEAR_BLK)

    # Ownership mini-table + facts row
    fx = cx; fy = purpose_card_y + Inches(0.85)
    _add_text(s, fx, fy, cw, Inches(0.3),
              "### Ownership",
              size=11, bold=True, color=NEAR_BLK)
    facts_line = (
        "core-banking   ·   tier 0   ·   #core-banking-oncall   ·   cobol-gnucobol   ·   commit c4c9ce66"
    )
    _add_text(s, fx, fy + Inches(0.3), cw, Inches(0.3),
              facts_line, size=10, color=DARK_GREY, font=FONT_MONO)

    # Provides / Depends strip
    pv_y = fy + Inches(0.75)
    _add_text(s, fx, pv_y, cw, Inches(0.3),
              "### What this application provides", size=11, bold=True, color=NEAR_BLK)
    _add_text(s, fx, pv_y + Inches(0.3), cw, Inches(0.3),
              "owned schema  ledger   ·   writes  ledger.balances   ·   writes file  /shared/reports/LEDGER.RPT",
              size=10, color=DARK_GREY, font=FONT_MONO)

    dp_y = pv_y + Inches(0.7)
    _add_text(s, fx, dp_y, cw, Inches(0.3),
              "### What this application depends on", size=11, bold=True, color=NEAR_BLK)
    _add_text(s, fx, dp_y + Inches(0.3), cw, Inches(0.3),
              "reads file  /shared/postings/POSTINGS.DAT",
              size=10, color=DARK_GREY, font=FONT_MONO)

    # Impact badge
    im_y = dp_y + Inches(0.72)
    _add_rect(s, fx, im_y, cw, Inches(0.3), fill=RGBColor(0xFE, 0xF3, 0xC7))
    _add_text(s, fx + Inches(0.15), im_y + Inches(0.04), cw - Inches(0.3), Inches(0.3),
              "IMPACT  ·  changes here may affect 2 apps:  account-service, analytics-service  ·  "
              "depends on 0 upstream",
              size=10, bold=True, color=RGBColor(0x92, 0x40, 0x0E))

    # Right rail — key properties (compact)
    rx = Inches(9.15); ry = Inches(2.85); rw = Inches(3.55)
    _add_text(s, rx, ry, rw, Inches(0.4),
              "Key properties", size=11, bold=True, color=UBS_RED)
    _add_bullets(s, rx, ry + Inches(0.45), rw, Inches(3.9), [
        ("Per-app. ", "One MR in the app's own repo."),
        ("Reviewed by owner. ", "Not a cross-team propagation — that's slide 10."),
        ("Non-destructive. ", "Only the marker-delimited block is touched."),
        ("Never direct-commit. ", "Always a draft MR."),
        ("LLM pluggable. ", "Mock · Gemini · Azure OpenAI. One env var."),
        ("Grounded. ", "Prompt feeds only graph evidence — no hallucinated contracts."),
    ], size=10, line_spacing=1.35)

    # Bottom caption under the preview
    _add_text(s, card_x, Inches(6.95), card_w, Inches(0.25),
              "Actual Gemini output from the POC, quoted verbatim from "
              "banking/legacy-ledger autodoc MR.",
              size=8, italic=True, color=MID_GREY)

    _notes(s, "Anchor phrasing: 'If legacy-ledger ever wondered what depends on them — "
              "the answer is now in their own README, automatically. No hand-maintenance, no stale "
              "Confluence pages.' The one-liner Purpose quote is REAL — Gemini 2.5 Flash, grounded "
              "on the graph's own evidence (writes ledger.balances, writes LEDGER.RPT, reads POSTINGS.DAT). "
              "Key contrast: this MR is in legacy-ledger's own repo, for core-banking to review — NOT a "
              "PR to the 2 affected downstream apps (that's slide 10's sticky-tag flow).")
    return s


def build_proof(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "Proof of concept — what's running today",
                 eyebrow="Evidence",
                 subtitle="All five loops live in a local environment. Ten fixture banking apps "
                          "+ one auto-discovered mid-demo. Real enterprise language mix.")

    # left: what we built
    left_x = Inches(0.6); left_y = Inches(3.0); left_w = Inches(5.8)
    _add_text(s, left_x, left_y, left_w, Inches(0.4),
              "Polyglot fixture covering UBS's real language mix",
              size=12, bold=True, color=UBS_RED)
    _add_bullets(s, left_x, left_y + Inches(0.5), left_w, Inches(3.7), [
        ("Java/Spring Boot · ", "account, transaction, customer, notification services"),
        ("Python/FastAPI · ",   "analytics, fraud-detection"),
        ("Node/Express · ",     "web-portal-bff, mobile-bff"),
        ("COBOL · ",            "legacy-ledger (postings batch + ledger report)"),
        ("Hybrid platform · ",  "5 on-prem apps + 4 azure apps + 1 library, tagged via CMDB"),
        ("Contracts · ",        "OpenAPI, AsyncAPI/Kafka, Flyway SQL, shared file feeds"),
        ("LLM · ",              "Gemini (local) + Azure OpenAI (enterprise), pluggable adapter"),
    ], size=12, line_spacing=1.4)

    # right: headline metrics
    right_x = Inches(7.0); right_y = Inches(3.0); right_w = Inches(5.7)
    _add_text(s, right_x, right_y, right_w, Inches(0.4),
              "POC numbers", size=12, bold=True, color=UBS_RED)

    metrics = [
        ("59",  "graph nodes · 81 edges across 6 coupling types incl. FILE"),
        ("11",  "REST call edges — exact method + path via regex + tree-sitter AST"),
        ("11",  "break kinds — endpoint/topic/table/file/library, required-vs-optional aware"),
        ("3",   "languages parsed at AST level — Java · Python · JS/TS"),
        ("60s", "reconciler interval — new GitLab projects auto-discovered + hooked"),
        ("6",   "autodoc MRs opened via Gemini in one demo session (click-driven)"),
    ]
    for i, (n, desc) in enumerate(metrics):
        my = right_y + Inches(0.5 + i * 0.52)
        _add_rect(s, right_x, my, Inches(0.95), Inches(0.45), fill=UBS_RED)
        _add_text(s, right_x, my + Inches(0.05), Inches(0.95), Inches(0.4),
                  n, size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(s, right_x + Inches(1.15), my + Inches(0.1), right_w - Inches(1.25), Inches(0.4),
                  desc, size=10, color=DARK_GREY)

    _add_text(s, Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.3),
              "Full live demo available on request — 15 min end-to-end across all five loops.",
              size=10, italic=True, color=MID_GREY)

    _notes(s, "Key reframing: 'Most of H1 and a meaningful slice of H2 are already working. "
              "Pilot money is NOT to prove the idea — that's done. It's to operationalize: "
              "bot account, security review, multi-group rollout, scheduled autodoc loop, "
              "rollout to domain #1.'")
    return s


def build_roadmap(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "Roadmap — four horizons", eyebrow="Where this goes",
                 subtitle="H1 and most of H2 are already running in the POC. Pilot funds "
                          "operationalization + first domain. H3–H4 follow once visibility is "
                          "embedded.")

    horizons = [
        ("H1", "POC ✓ · pilot 0–3 mo",
         "Cross-app visibility + MR impact bot + auto-discovery + hybrid platform",
         "Canvas live. MR impact bot live. Sticky impact tags live. Hybrid Azure ↔ on-prem "
         "callouts live. Reconciler auto-discovers new GitLab projects, installs webhooks, "
         "detects renames and archival. Pilot = production hardening + first domain rollout.",
         UBS_RED),
        ("H2", "POC started · pilot 3–6 mo",
         "Auto-documentation of legacy code (LLM)",
         "Per-app draft MR with a Sherlock-managed README section; "
         "Gemini + Azure OpenAI adapters already wired; one-click from canvas. "
         "Pilot work: scheduled loop, priority scoring, hierarchical summarization, "
         "COBOL-tuned prompts.",
         RGBColor(0xF5, 0x9E, 0x0B)),
        ("H3", "6–12 mo",
         "SLA dashboards + governance",
         "Mean-time-to-clear impact tag per team. Optional hard-gate for catastrophic "
         "breaks. Domain-level architectural-debt score. Backstage plugin.",
         RGBColor(0x8B, 0x5C, 0xF6)),
        ("H4", "12+ mo",
         "Runtime-signal integration",
         "Cross-validate static edges against Dynatrace traces. Flag dead edges (never "
         "fire). Flag missing edges (fires in prod, not in code). Historical graph diffs.",
         RGBColor(0x10, 0xB9, 0x81)),
    ]

    y = Inches(2.9); row_h = Inches(1.05); gap = Inches(0.08)
    for i, (code, dur, headline, body, color) in enumerate(horizons):
        ry = y + i * (row_h + gap)
        _add_rect(s, Inches(0.6), ry, Inches(0.9), row_h, fill=color)
        _add_text(s, Inches(0.6), ry + Inches(0.12), Inches(0.9), Inches(0.35),
                  code, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(s, Inches(0.6), ry + Inches(0.56), Inches(0.9), Inches(0.35),
                  dur, size=8, color=WHITE, align=PP_ALIGN.CENTER)
        _add_rect(s, Inches(1.55), ry, Inches(11.2), row_h, fill=WHITE, line=LINE_GREY)
        _add_text(s, Inches(1.75), ry + Inches(0.15), Inches(10.8), Inches(0.4),
                  headline, size=14, bold=True, color=NEAR_BLK)
        _add_text(s, Inches(1.75), ry + Inches(0.55), Inches(10.8), Inches(0.5),
                  body, size=10, color=DARK_GREY)

    _notes(s, "H2 is the 'code documentation' pillar from the original vision. Note: "
              "'Auto-doc is a natural extension of the same graph — we already understand "
              "every repo's structure; generating docs becomes a scheduled job using the "
              "LLM provider you already have in Azure.'")
    return s


def build_pilot_plan(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "Pilot plan — 3 months",
                 eyebrow="How we move from POC to production",
                 subtitle="Most of the build is done in the POC. These three months turn it "
                          "into a platform-embedded, auditable, domain-ready service.")

    months = [
        ("Month 1", "Production-harden. Embed.",
         ["Deploy Sherlock alongside the existing GitLab CI stack (shared infra).",
          "Swap PAT for a GitLab bot account / GitLab App with group scopes.",
          "Security review + data-residency sign-off (code-only, no PII — A4 appendix).",
          "Wire CMDB read-only + Slack/Teams notifications to on-call channels.",
          "Ingest one pilot domain (3–10 apps) — reconciler auto-onboards new repos."]),
        ("Month 2", "Enable the MR bot + autodoc.",
         ["Turn on MR impact bot for pilot domain in opt-in mode.",
          "Measure break detection, author response, false-positive rate.",
          "Turn on autodoc in draft mode (one-click); gate scheduled loop behind opt-in.",
          "Iterate parsers for any domain-specific idioms (in-house DB libs, schedulers)."]),
        ("Month 3", "Prove value. Decide scale.",
         ["Enable sticky impact tags for pilot domain.",
          "Onboard a second domain to test multi-group + cross-team notification.",
          "Present outcome to CTO office — MTTD reduction, engagement, coverage.",
          "Decision: enterprise rollout plan for H2 autodoc scheduler + H3 SLA board."]),
    ]
    y = Inches(2.9); w = Inches(4.0); gap = Inches(0.15); x0 = Inches(0.6); h = Inches(4.0)
    for i, (month, lead, bullets) in enumerate(months):
        x = x0 + (w + gap) * i
        _add_rect(s, x, y, w, h, fill=WHITE, line=LINE_GREY)
        _add_rect(s, x, y, w, Inches(0.8), fill=NEAR_BLK)
        _add_text(s, x + Inches(0.25), y + Inches(0.1), w - Inches(0.5), Inches(0.4),
                  month.upper(), size=10, bold=True, color=UBS_RED)
        _add_text(s, x + Inches(0.25), y + Inches(0.38), w - Inches(0.5), Inches(0.4),
                  lead, size=13, bold=True, color=WHITE)
        _add_bullets(s, x + Inches(0.25), y + Inches(1.05), w - Inches(0.5), h - Inches(1.2),
                     bullets, size=10, line_spacing=1.3)

    _notes(s, "Key framing vs. the previous pilot plan: 'Month 1 isn't about building — that's "
              "already done. Month 1 is about security, governance, and embedding with the "
              "platform org. Month 3 closes with a go/no-go for H2's scheduled autodoc loop + "
              "H3 governance work.'")
    return s


def build_ask(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    # dark background for impact
    _add_rect(s, Inches(0), Inches(0.8), prs.slide_width, prs.slide_height - Inches(1.2), fill=NEAR_BLK)
    _title_block(s, "The ask", eyebrow="What we need from you")

    asks = [
        ("3-MONTH FUNDED PILOT",
         "Approval to run the plan on the previous slide."),
        ("2–3 FTEs",
         "1 senior engineer · 1 engineer · 1 UX/platform (shared from platform org)."),
        ("PLATFORM EMBED",
         "Approval to deploy alongside the existing GitLab CI stack — shared infra, no new standalone service."),
        ("EXECUTIVE SPONSORSHIP",
         "One CTO-office sponsor to unblock cross-BU alignment when we onboard domain #2."),
    ]
    y = Inches(2.8); h = Inches(0.85); gap = Inches(0.2)
    for i, (head, body) in enumerate(asks):
        ry = y + i * (h + gap)
        _add_rect(s, Inches(0.8), ry, Inches(0.1), h, fill=UBS_RED)
        _add_text(s, Inches(1.1), ry + Inches(0.07), Inches(5.0), Inches(0.45),
                  head, size=14, bold=True, color=WHITE)
        _add_text(s, Inches(1.1), ry + Inches(0.45), Inches(11), Inches(0.4),
                  body, size=11, color=LINE_GREY)

    _add_text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.35),
              "What we'll deliver: one pilot domain with measurable MTTD reduction, "
              "a ready-to-scale platform for H2.",
              size=11, italic=True, color=LINE_GREY, align=PP_ALIGN.CENTER)

    _notes(s, "Stop talking. Let the CTO respond. Have answers ready for: 'Why now?', "
              "'Why your team?', 'Why not Backstage?', 'What's the total 12-month cost?', "
              "'What's the kill-switch if the bot produces noise?'")
    return s


def build_metrics(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _title_block(s, "How we'll know it's working", eyebrow="Success metrics",
                 subtitle="We measure the platform, not the people. Teams are informed, not policed.")

    metrics = [
        ("MTTD for cross-app breaks",
         "Baseline",  "1–5 days",
         "Target (pilot)", "< 1 hour (at MR time)", UBS_RED),
        ("MR impact-comment engagement",
         "Target", ">80% read by MR author",
         "Secondary", ">60% elicit a coordination action", BLACK),
        ("Sticky-issue cycle time",
         "Report", "Median `pending` → `fixed`, per team",
         "SLA tier", "Tier 0 apps: <24h; Tier 1: <72h", BLACK),
        ("Coverage",
         "Target", ">90% of pilot-domain active repos indexed",
         "Freshness", "Graph reflects main within 24h of push", BLACK),
    ]
    y = Inches(2.9); w = Inches(6.05); gap = Inches(0.15); h = Inches(1.85)
    for i, (name, l1a, l1b, l2a, l2b, accent) in enumerate(metrics):
        x = Inches(0.6) + (i % 2) * (w + gap)
        yy = y + (i // 2) * (h + gap)
        _add_rect(s, x, yy, w, h, fill=WHITE, line=LINE_GREY)
        _add_rect(s, x, yy, Inches(0.08), h, fill=accent)
        _add_text(s, x + Inches(0.25), yy + Inches(0.15), w - Inches(0.5), Inches(0.4),
                  name, size=13, bold=True, color=NEAR_BLK)
        _add_text(s, x + Inches(0.25), yy + Inches(0.65), Inches(1.4), Inches(0.3),
                  l1a.upper(), size=9, bold=True, color=accent)
        _add_text(s, x + Inches(1.8), yy + Inches(0.62), w - Inches(2.0), Inches(0.4),
                  l1b, size=12, color=DARK_GREY)
        _add_text(s, x + Inches(0.25), yy + Inches(1.15), Inches(1.4), Inches(0.3),
                  l2a.upper(), size=9, bold=True, color=accent)
        _add_text(s, x + Inches(1.8), yy + Inches(1.12), w - Inches(2.0), Inches(0.4),
                  l2b, size=12, color=DARK_GREY)

    _notes(s, "These are the numbers we'll report at month 3. If they don't move, the "
              "pilot has failed and we say so. Transparency builds trust for H2.")
    return s


def build_close(prs, idx, total):
    s = _blank_slide(prs); _header(s); _footer(s, idx, total)
    _add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill=NEAR_BLK)
    _add_rect(s, Inches(0), Inches(0), prs.slide_width, Inches(0.15), fill=UBS_RED)

    _add_text(s, Inches(0.6), Inches(1.8), Inches(12), Inches(0.5),
              "CLOSING", size=11, bold=True, color=UBS_RED)
    _add_text(s, Inches(0.6), Inches(2.3), Inches(12), Inches(1.5),
              "Every change has a blast radius.", size=40, bold=True, color=WHITE)
    _add_text(s, Inches(0.6), Inches(3.5), Inches(12), Inches(1.5),
              "Sherlock makes it visible — at the moment a developer hits Merge.",
              size=22, color=LINE_GREY)
    _add_rect(s, Inches(0.6), Inches(5.0), Inches(1.6), Inches(0.06), fill=UBS_RED)
    _add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.5),
              "Thank you. Questions?", size=18, italic=True, color=WHITE)
    _add_text(s, Inches(0.6), Inches(6.9), Inches(12), Inches(0.3),
              "Appendix follows — architecture, schema, security, roadmap details.",
              size=9, color=MID_GREY)
    _notes(s, "Hand over. Keep the appendix in your back pocket for Q&A.")
    return s


# ---------- appendix slides ---------------------------------------------------

def build_appendix_divider(prs, idx, total):
    s = _blank_slide(prs)
    _add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill=BG_GREY)
    _add_rect(s, Inches(0.6), Inches(3.4), Inches(1.6), Inches(0.08), fill=UBS_RED)
    _add_text(s, Inches(0.6), Inches(2.5), Inches(12), Inches(0.5),
              "APPENDIX", size=12, bold=True, color=UBS_RED)
    _add_text(s, Inches(0.6), Inches(3.6), Inches(12), Inches(1.0),
              "Deeper dives for Q&A",
              size=36, bold=True, color=NEAR_BLK)
    _add_text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.5),
              "Architecture · Graph schema · GitLab integration · Security · LLM strategy · Positioning",
              size=14, color=MID_GREY)
    _notes(s, "Only surface appendix slides if the CTO asks a question that one covers.")
    return s


def build_a_architecture_deep(prs, idx, total):
    s = _blank_slide(prs); _header(s, "APPENDIX"); _footer(s, idx, total, "APPENDIX")
    _title_block(s, "Architecture — deeper view", eyebrow="Appendix · A1")

    _add_bullets(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(4.0), [
        ("Ingest (FastAPI). ", "Webhook receiver with marker-idempotent MR-note and issue upsert; "
                                "manual /scan-all for full-estate refresh."),
        ("Analyzers. ",         "Regex + language-aware heuristics: Java (FileInputStream, RestClient.uri, "
                                "Flyway YAML), Python (httpx, psycopg, pathlib open modes), "
                                "Node (axios / fs), COBOL (SELECT/ASSIGN + OPEN verbs), "
                                "OpenAPI / AsyncAPI / Maven-pom / requirements.txt / package.json."),
        ("Graph DB. ",          "Neo4j Community 5.x. Nodes indexed by natural keys. "
                                "Re-scan is idempotent per app — outgoing edges replaced atomically."),
        ("Impact Engine. ",     "Pure Cypher. Diff two AnalysisResult snapshots; route by kind; "
                                "enrich impacted apps via CMDB."),
        ("Notifier. ",           "GitLab API: merge-request notes (marker-based upsert), project issues "
                                "(label::pending → label::fixed), application settings (allow-local-requests)."),
        ("Canvas. ",             "Cytoscape.js single-page. Dagre rankDir=BT (upstream top, downstream bottom). "
                                "App view + Contract view."),
    ], size=11, line_spacing=1.35)

    _notes(s, "Used when asked 'how does it scale?' or 'what's the stack?'")
    return s


def build_a_coverage_detail(prs, idx, total):
    s = _blank_slide(prs); _header(s, "APPENDIX"); _footer(s, idx, total, "APPENDIX")
    _title_block(s, "Detection coverage — detail & gaps",
                 eyebrow="Appendix · A2")

    _add_text(s, Inches(0.6), Inches(2.9), Inches(12), Inches(0.4),
              "Covered today", size=13, bold=True, color=UBS_RED)
    _add_bullets(s, Inches(0.6), Inches(3.3), Inches(12.1), Inches(1.9), [
        "REST endpoints exposed (OpenAPI) + called — regex AND tree-sitter AST for Java · Python · JS/TS",
        "Required-vs-optional payload diffing — required-field changes flagged BREAKING, optional additions land as INFO",
        "Kafka topics published/consumed (AsyncAPI + in-code KafkaTemplate / confluent-kafka), payload precision identical to REST",
        "Postgres schemas / tables — DDL (Flyway) + DML via SQL-in-code heuristics",
        "Shared file feeds on /shared/, /mnt/feeds/, /inbound/, /outbound/ — paired COBOL SELECT/ASSIGN",
        "Shared libraries via pom.xml / requirements.txt / pyproject / package.json",
        "Hybrid platform tagging — CMDB-driven `platform` field merged into the graph; cross-boundary callouts at MR time",
    ], size=11)

    _add_text(s, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
              "Known gaps (roadmap)", size=13, bold=True, color=UBS_RED)
    _add_bullets(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.5), [
        "gRPC contracts (.proto) and GraphQL federation schemas",
        "SFTP / FTP transfer idioms and explicit file-watcher patterns",
        "In-house message-bus libraries that don't emit AsyncAPI",
        "Batch-job schedulers (Control-M) — currently only detects schedule annotations in code",
        "Polyglot gRPC code-gen where the contract is generated, not committed",
    ], size=11)

    _notes(s, "Use to show we know what we don't know. Every gap has an obvious next "
              "analyzer; none require rebuilding the foundation.")
    return s


def build_a_schema(prs, idx, total):
    s = _blank_slide(prs); _header(s, "APPENDIX"); _footer(s, idx, total, "APPENDIX")
    _title_block(s, "Graph schema", eyebrow="Appendix · A3")

    _add_text(s, Inches(0.6), Inches(2.9), Inches(12), Inches(0.4),
              "Node types", size=13, bold=True, color=UBS_RED)
    _add_text(s, Inches(0.6), Inches(3.3), Inches(12.1), Inches(1.3),
              "Application   ·   Endpoint   ·   Topic   ·   DBSchema   ·   DBTable   ·   "
              "FileFeed   ·   Library",
              size=13, color=DARK_GREY, font=FONT_MONO)

    _add_text(s, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4),
              "Edge types (12 total)", size=13, bold=True, color=UBS_RED)
    _add_bullets(s, Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.2), [
        "EXPOSES (App→Endpoint)   ·   CALLS (App→Endpoint)",
        "PUBLISHES (App→Topic)   ·   CONSUMES (App→Topic)",
        "OWNS_SCHEMA (App→DBSchema)   ·   CONTAINS_TABLE (DBSchema→DBTable)",
        "READS_TABLE / WRITES_TABLE (App→DBTable)",
        "READS_FILE / WRITES_FILE (App→FileFeed)",
        "DEPENDS_ON_LIB (App→Library)   ·   PUBLISHES_LIB (App→Library)",
    ], size=11, line_spacing=1.35)

    _notes(s, "The schema is deliberately narrow — every edge type corresponds to a "
              "real-world break kind we can detect.")
    return s


def build_a_security(prs, idx, total):
    s = _blank_slide(prs); _header(s, "APPENDIX"); _footer(s, idx, total, "APPENDIX")
    _title_block(s, "Security, data residency, PII",
                 eyebrow="Appendix · A4")

    _add_bullets(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(3.8), [
        ("No PII in the graph. ", "Sherlock reads source code — not runtime data. The graph "
                                  "contains app names, endpoint paths, topic names, table FQNs, "
                                  "file paths, team names."),
        ("All in-tenant. ",       "Graph DB, analyzers, and UI run inside the UBS perimeter. "
                                  "No code leaves."),
        ("Phase-2 LLM: in-tenant only. ", "Auto-doc uses Azure OpenAI in our own tenant (GPT-5.1) "
                                         "— no external API calls. Local dev may use a different "
                                         "provider (abstraction layer)."),
        ("Access control on canvas (future). ", "Canvas respects GitLab project visibility. "
                                               "Team/domain filtering before public rollout."),
        ("Supply chain. ",        "Open-source deps pinned; Neo4j Community; no commercial "
                                  "licenses required for the platform itself."),
        ("Auditability. ",        "Every MR comment and sticky-issue action logged; "
                                  "commit SHAs stamped on each Application node."),
    ], size=12, line_spacing=1.4)

    _notes(s, "Key talking point for any CISO/risk review: 'Code, not data.'")
    return s


def build_a_llm(prs, idx, total):
    s = _blank_slide(prs); _header(s, "APPENDIX"); _footer(s, idx, total, "APPENDIX")
    _title_block(s, "H2 — Auto-documentation with LLM",
                 eyebrow="Appendix · A5",
                 subtitle="Not in H1 pilot scope. Activated in month 4 subject to pilot outcome.")

    _add_bullets(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(3.8), [
        ("Goal. ",     "Generate README-level docs for legacy code — especially COBOL and "
                       "long-tenured Java — pushed as MRs for human review. Never direct-commit."),
        ("Triggers. ", "Scheduled per-repo, prioritized by freshness (files untouched >12 mo) "
                       "and docstring density (low = higher priority)."),
        ("Provider. ", "Pluggable. Enterprise: Azure OpenAI GPT-5.1 in UBS tenant. "
                       "Local dev: configurable."),
        ("Cost control. ", "Hierarchical summarization (function → file → module → service). "
                            "Aggressive prompt caching. Cost per doc should be < one engineer-hour."),
        ("Quality gate. ", "Generated docs land as draft MR; owning team reviews; merge is a "
                            "human decision."),
        ("Audit. ",    "Every generated doc tagged in Git with provider + model + prompt hash."),
    ], size=12, line_spacing=1.4)

    _notes(s, "LLM talking point: 'The expensive part is understanding the code. Sherlock "
              "already has that model — doc generation is a scheduled job on top.'")
    return s


def build_a_integrations(prs, idx, total):
    s = _blank_slide(prs); _header(s, "APPENDIX"); _footer(s, idx, total, "APPENDIX")
    _title_block(s, "Integration futures", eyebrow="Appendix · A6")

    _add_bullets(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(3.8), [
        ("CMDB (read-only). ",   "Pull team / tier / on-call / platform. Source of truth "
                                 "for ownership. Sherlock writes nothing back to CMDB."),
        ("Backstage. ",          "Expose Sherlock's app graph as a Backstage plugin — "
                                 "embeds the canvas in the existing dev-portal."),
        ("APM (Dynatrace). ",    "Cross-validate runtime traces against static edges. "
                                 "Flag unused edges (static-only) and surprise edges "
                                 "(runtime-only) — both are signals."),
        ("API / AI gateway adapter. ", "Read gateway routing tables to unravel calls "
                                       "that look generic at the source ('GET /v1/...') but "
                                       "fan out to many backend services."),
        ("Slack / Teams. ",      "Route impact notifications to on-call channels named in CMDB."),
        ("GitLab app / bot. ",   "Graduate from PAT to a proper GitLab app; per-group bot account; "
                                 "fine-grained scopes."),
        ("Control-M / schedulers. ", "Detect scheduled-job triggers that cross boundaries "
                                      "(read from shared fileserver, etc.)."),
    ], size=11, line_spacing=1.35)

    _notes(s, "Used when asked 'How does this fit with our existing tooling?' — answer: "
              "it plugs in, doesn't replace.")
    return s


def build_a_positioning(prs, idx, total):
    s = _blank_slide(prs); _header(s, "APPENDIX"); _footer(s, idx, total, "APPENDIX")
    _title_block(s, "Positioning vs. adjacent tools",
                 eyebrow="Appendix · A7")

    rows = [
        ("Backstage",            "Ownership + docs portal.",
                                 "Sherlock supplies the dependency graph Backstage lacks."),
        ("Dependabot / Renovate","Third-party library dep updates.",
                                 "Different scope — they do CVEs/versions, Sherlock does "
                                 "cross-app contract impact."),
        ("GitLab Dependency Scanning", "Security-focused.",
                                 "Security CVE signals ≠ cross-app behavioural impact."),
        ("APM (Dynatrace / AppDynamics)", "Runtime calls.",
                                 "Post-hoc, incomplete for batch/files; Sherlock is "
                                 "pre-merge, contract-based."),
        ("Code search (Sourcegraph)", "Reference finder.",
                                 "Sourcegraph tells you where a symbol appears; "
                                 "Sherlock tells you what breaks when it changes."),
        ("Sonar / Checkmarx",    "Code quality + security.",
                                 "Intra-app signals; silent on inter-app surfaces."),
    ]
    y = Inches(2.9)
    _add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.45), fill=NEAR_BLK)
    _add_text(s, Inches(0.8), y + Inches(0.1), Inches(3.5), Inches(0.3),
              "TOOL", size=10, bold=True, color=WHITE)
    _add_text(s, Inches(4.4), y + Inches(0.1), Inches(3.5), Inches(0.3),
              "WHAT IT DOES", size=10, bold=True, color=WHITE)
    _add_text(s, Inches(7.8), y + Inches(0.1), Inches(5.0), Inches(0.3),
              "WHERE SHERLOCK IS DIFFERENT", size=10, bold=True, color=WHITE)
    for i, (tool, what, diff) in enumerate(rows):
        ry = y + Inches(0.5 + i * 0.55)
        _add_rect(s, Inches(0.6), ry, Inches(12.1), Inches(0.5),
                  fill=BG_GREY if i % 2 == 0 else WHITE, line=LINE_GREY)
        _add_text(s, Inches(0.8), ry + Inches(0.1), Inches(3.5), Inches(0.3),
                  tool, size=10, bold=True, color=NEAR_BLK)
        _add_text(s, Inches(4.4), ry + Inches(0.1), Inches(3.3), Inches(0.3),
                  what, size=10, color=DARK_GREY)
        _add_text(s, Inches(7.8), ry + Inches(0.1), Inches(5.0), Inches(0.3),
                  diff, size=10, color=DARK_GREY)

    _notes(s, "Expected question: 'Don't we already have X for this?' Walk through the row "
              "for the tool they name; the honest answer is usually 'X does Y, Sherlock does "
              "cross-app-impact — they complement, not compete.'")
    return s


def build_a_operator_features(prs, idx, total):
    s = _blank_slide(prs); _header(s, "APPENDIX"); _footer(s, idx, total, "APPENDIX")
    _title_block(
        s,
        "Operator features — auto-discovery, archival, rename",
        eyebrow="Appendix · A8",
        subtitle="Platform-team capabilities that make Sherlock self-maintaining at "
                 "thousands-of-repos scale. All built in the POC.",
    )

    _add_bullets(s, Inches(0.6), Inches(3.1), Inches(12.1), Inches(3.9), [
        ("Auto-discovery (60s reconciler). ",
         "Background task lists every project in every configured group. New projects → "
         "webhook auto-installed + initial scan. Webhook deleted by a team? Reinstalled "
         "on the next tick. Metadata refresh even when code hasn't changed."),
        ("Multi-group support. ",
         "GITLAB_GROUPS env var accepts a comma-separated list (payments, wealth, "
         "core-banking, …). Reconciler iterates each; one Sherlock covers the whole estate."),
        ("Parallel scan. ",
         "Bounded ThreadPoolExecutor (default 4 workers) scans newly-discovered repos "
         "concurrently. Scales with CPU, not with group size."),
        ("Archival. ",
         "Projects removed from every configured group are marked archived in the graph "
         "— Application node + edges preserved for audit. `impact::pending` issues still "
         "close normally."),
        ("Rename detection. ",
         "Stable project-ID tracking. A path rename in GitLab is detected as a rename "
         "(not a delete + new), old node archived with renamed_to pointing at the new one."),
        ("Observability. ",
         "/api/reconciler/status exposes last run, last N runs, running state. "
         "Every run emits a structured log line for the platform team's dashboard."),
    ], size=11, line_spacing=1.35)

    _notes(s, "Surface this when asked 'what does the platform team operate?' or "
              "'what happens when our 1,000 repos become 10,000?' — answer: the reconciler "
              "does the onboarding; there's no manual list to maintain.")
    return s


def build_a_whynow(prs, idx, total):
    s = _blank_slide(prs); _header(s, "APPENDIX"); _footer(s, idx, total, "APPENDIX")
    _title_block(s, "Why now. Why UBS. Why us.",
                 eyebrow="Appendix · A8")

    _add_bullets(s, Inches(0.6), Inches(3.0), Inches(12.1), Inches(3.8), [
        ("GitLab is where MRs happen. ", "Impact belongs on the MR — not in a portal nobody "
                                          "opens. Sherlock lives where developers already work."),
        ("17 domains = invisible cross-boundary coupling. ",
         "Every large bank has this; our scale makes it expensive. The "
         "cost of not knowing who depends on you grows exponentially with teams."),
        ("COBOL coverage is ready. ",   "Few tools in this space speak COBOL. Our core-banking "
                                        "estate deserves to be in the graph, not a footnote."),
        ("LLM tooling has matured. ",   "Doc generation at scale is finally feasible with "
                                        "prompt caching + hierarchical summarization."),
        ("Proof already done. ",        "POC runs on a laptop, exercises every language we "
                                        "care about, demonstrates the MR bot end-to-end. "
                                        "Pilot money goes to scaling safely, not to proving "
                                        "the idea."),
    ], size=12, line_spacing=1.4)

    _notes(s, "Closing appendix slide — use when asked 'why is this the right moment?' or "
              "'why your team?'")
    return s


# ---------- main --------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    main_builders = [
        build_title_slide,      # 1
        build_exec_summary,     # 2
        build_problem,          # 3
        build_gap,              # 4
        build_intro_sherlock,   # 5
        build_architecture,     # 6
        build_coverage,         # 7
        build_mr_moment,        # 8  — MR comment for BREAKING changes (for author)
        build_canvas,           # 9
        build_sticky_lifecycle, # 10 — impact::pending issues in DOWNSTREAM repos (for affected teams)
        build_hybrid_platform,  # 11 — Azure / on-prem boundary callout — single graph, single workflow
        build_autodoc,          # 12 — per-app README MR in THAT APP's own repo (for owning team)
        build_proof,            # 13
        build_roadmap,          # 14
        build_pilot_plan,       # 15
        build_ask,              # 16
        build_metrics,          # 17
        build_close,            # 18
    ]
    appx_builders = [
        build_appendix_divider,
        build_a_architecture_deep,
        build_a_coverage_detail,
        build_a_schema,
        build_a_security,
        build_a_llm,
        build_a_integrations,
        build_a_positioning,
        build_a_operator_features,
        build_a_whynow,
    ]
    total = len(main_builders) + len(appx_builders)
    for i, fn in enumerate(main_builders, start=1):
        fn(prs, i, total)
    for i, fn in enumerate(appx_builders, start=len(main_builders) + 1):
        fn(prs, i, total)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"wrote {OUT_PATH}  ({total} slides)")


if __name__ == "__main__":
    main()
